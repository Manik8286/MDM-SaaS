"""
Generate MDM SaaS demo PowerPoint presentation.
Run: python3 scripts/gen_pptx.py
Output: MDM_SaaS_Demo.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ─────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1D, 0x4E, 0xD8)   # primary
DARK       = RGBColor(0x0F, 0x17, 0x2A)   # dark bg
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY   = RGBColor(0x64, 0x74, 0x8B)
GREEN      = RGBColor(0x05, 0x96, 0x69)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, space_before=6):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def dark_slide(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), BLUE)
    add_text(slide, title, Inches(0.4), Inches(0.15), Inches(10), Inches(0.6),
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), Inches(12), Inches(0.35),
                 size=14, bold=False, color=RGBColor(0xBA, 0xD8, 0xFF))

def bullet_card(slide, l, t, w, h, title, bullets, title_color=BLUE, bg=LIGHT_GRAY):
    add_rect(slide, l, t, w, h, bg)
    # title
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.1), w - Inches(0.3), Inches(0.4))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = title_color
    # bullets
    txb2 = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.48), w - Inches(0.3), h - Inches(0.55))
    txb2.word_wrap = True
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p2 = tf2.paragraphs[0]
            first = False
        else:
            p2 = tf2.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = f"• {b}"
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
dark_slide(slide)

# accent stripe
add_rect(slide, 0, Inches(2.8), Inches(0.08), Inches(2.0), BLUE)

add_text(slide, "MDM SaaS", Inches(0.3), Inches(0.7), Inches(12), Inches(1.2),
         size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "Multi-Tenant Apple Device Management Platform",
         Inches(0.3), Inches(1.85), Inches(11), Inches(0.7),
         size=26, bold=False, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.LEFT)

add_text(slide, "Product Demo", Inches(0.3), Inches(2.85), Inches(4), Inches(0.45),
         size=16, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

features = [
    "Enrollment & Device Management",
    "Platform SSO with Microsoft Entra ID",
    "Software Distribution & Self-Service Portal",
    "USB Policy, Gatekeeper & Compliance",
    "JIT Admin Access with Auto-Revoke",
]
y = Inches(3.5)
for f in features:
    add_text(slide, f"  ✦  {f}", Inches(0.35), y, Inches(8), Inches(0.38),
             size=15, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(0.42)

add_text(slide, "Built with FastAPI · PostgreSQL · Apple APNs · Next.js",
         Inches(0.3), Inches(6.9), Inches(10), Inches(0.4),
         size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Solution
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "The Problem We Solve",
           "Managing macOS devices at scale — without the complexity of enterprise MDM tools")

# Problem box
add_rect(slide, Inches(0.3), Inches(1.3), Inches(5.9), Inches(5.6), WHITE)
add_text(slide, "❌  Traditional Pain Points", Inches(0.5), Inches(1.45), Inches(5.5), Inches(0.45),
         size=16, bold=True, color=ORANGE)
problems = [
    ("High cost", "Legacy MDM tools (Jamf, Kandji) cost $6–15/device/month"),
    ("Complex setup", "Requires dedicated IT staff and weeks of onboarding"),
    ("No PSSO", "Most tools don't support macOS Platform SSO with Entra ID"),
    ("No agent control", "Can't run custom scripts or manage admin access remotely"),
    ("Siloed workflows", "Software requests, access, compliance in separate tools"),
]
y = Inches(2.0)
for title, desc in problems:
    add_text(slide, f"• {title}", Inches(0.55), y, Inches(5.5), Inches(0.3),
             size=14, bold=True, color=RGBColor(0x1E, 0x29, 0x3B))
    add_text(slide, f"  {desc}", Inches(0.55), y + Inches(0.28), Inches(5.5), Inches(0.3),
             size=12, color=MID_GRAY)
    y += Inches(0.75)

# Solution box
add_rect(slide, Inches(6.6), Inches(1.3), Inches(6.4), Inches(5.6), DARK)
add_text(slide, "✅  Our Solution", Inches(6.8), Inches(1.45), Inches(6.0), Inches(0.45),
         size=16, bold=True, color=GREEN)
solutions = [
    ("Affordable SaaS", "Multi-tenant — one server, many customers"),
    ("5-minute setup", "Docker Compose → enroll device in minutes"),
    ("Native PSSO", "Built-in Entra ID Platform SSO with Secure Enclave"),
    ("Management Agent", "Pure-bash agent — installs apps, elevates access"),
    ("Unified platform", "Device mgmt + software + compliance + access in one"),
]
y = Inches(2.0)
for title, desc in solutions:
    add_text(slide, f"• {title}", Inches(6.85), y, Inches(6.0), Inches(0.3),
             size=14, bold=True, color=WHITE)
    add_text(slide, f"  {desc}", Inches(6.85), y + Inches(0.28), Inches(6.0), Inches(0.3),
             size=12, color=RGBColor(0x94, 0xA3, 0xB8))
    y += Inches(0.75)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Platform Architecture", "How it all fits together")

components = [
    (Inches(0.3),  Inches(1.4), Inches(2.5), Inches(1.6), "Next.js Dashboard", BLUE,
     ["Admin web UI", "Device management", "Policy controls", "Audit logs"]),
    (Inches(3.2),  Inches(1.4), Inches(2.8), Inches(1.6), "FastAPI Backend", DARK,
     ["REST API /api/v1/", "MDM protocol /mdm/", "JWT + mTLS auth", "Async/await"]),
    (Inches(6.4),  Inches(1.4), Inches(2.8), Inches(1.6), "PostgreSQL", RGBColor(0x33, 0x69, 0xE8),
     ["Tenants & devices", "Commands queue", "Compliance results", "Audit trail"]),
    (Inches(9.6),  Inches(1.4), Inches(3.3), Inches(1.6), "Apple APNs", RGBColor(0x05, 0x96, 0x69),
     ["Wake-up push", "HTTP/2 provider API", "Per-device push", "Instant delivery"]),

    (Inches(0.3),  Inches(3.3), Inches(2.5), Inches(1.6), "macOS Device", RGBColor(0x7C, 0x3A, 0xED),
     ["MDM protocol", "Checkin + Connect", "Profile install", "Command execute"]),
    (Inches(3.2),  Inches(3.3), Inches(2.8), Inches(1.6), "Bash Agent", ORANGE,
     ["LaunchDaemon service", "Polls for jobs", "Installs software", "Reports results"]),
    (Inches(6.4),  Inches(3.3), Inches(2.8), Inches(1.6), "AWS SQS", RGBColor(0xD9, 0x77, 0x06),
     ["Command queue", "LocalStack (dev)", "Async processing", "Reliable delivery"]),
    (Inches(9.6),  Inches(3.3), Inches(3.3), Inches(1.6), "Self-Service Portal", RGBColor(0x0F, 0x76, 0x6E),
     ["Software catalog", "App requests", "Admin access req.", "Token-authenticated"]),
]

for l, t, w, h, title, color, items in components:
    add_rect(slide, l, t, w, h, color)
    add_text(slide, title, l + Inches(0.12), t + Inches(0.1), w - Inches(0.2), Inches(0.38),
             size=14, bold=True, color=WHITE)
    y2 = t + Inches(0.52)
    for item in items:
        add_text(slide, f"• {item}", l + Inches(0.12), y2, w - Inches(0.2), Inches(0.28),
                 size=11, color=RGBColor(0xE2, 0xE8, 0xF0))
        y2 += Inches(0.26)

# Flow description
add_rect(slide, Inches(0.3), Inches(5.2), Inches(12.7), Inches(1.8), WHITE)
add_text(slide, "Command Flow:", Inches(0.5), Inches(5.3), Inches(2.0), Inches(0.35),
         size=13, bold=True, color=DARK)
flow = "Admin action  →  FastAPI queues command in PostgreSQL  →  SQS message sent  →  " \
       "APNs wake-up push to device  →  Device calls /mdm/apple/connect  →  " \
       "Server returns command plist  →  Device executes  →  Result stored in DB"
add_text(slide, flow, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.1),
         size=12, color=MID_GRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Device Enrollment
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Device Enrollment", "Zero-touch enrollment with signed .mobileconfig profiles")

steps = [
    ("1", "Generate Token", BLUE,
     "Admin creates enrollment token in dashboard.\nChoose platform (macOS), reusable or single-use, expiry time."),
    ("2", "Share URL", PURPLE,
     "Enrollment URL shared with user.\nOpen in Safari on the Mac — no agent pre-installed needed."),
    ("3", "Install Profile", GREEN,
     "Signed .mobileconfig downloaded.\nmacOS prompts user to install the MDM enrollment profile."),
    ("4", "Authenticate", ORANGE,
     "Device sends Authenticate + TokenUpdate to /mdm/apple/checkin.\nServer stores APNs push token."),
    ("5", "Enrolled", RGBColor(0x05, 0x96, 0x69),
     "Device appears in dashboard as 'enrolled'.\nAPNs push, serial, OS version, model all captured."),
]

x = Inches(0.3)
for num, title, color, desc in steps:
    # circle
    circ = slide.shapes.add_shape(9, x, Inches(1.4), Inches(0.7), Inches(0.7))  # oval
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    txc = slide.shapes.add_textbox(x, Inches(1.44), Inches(0.7), Inches(0.55))
    txc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = txc.text_frame.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE

    add_rect(slide, x, Inches(2.25), Inches(2.4), Inches(4.6), WHITE)
    add_text(slide, title, x + Inches(0.1), Inches(2.35), Inches(2.2), Inches(0.4),
             size=14, bold=True, color=color)
    add_text(slide, desc, x + Inches(0.1), Inches(2.82), Inches(2.2), Inches(3.8),
             size=12, color=RGBColor(0x33, 0x41, 0x55), wrap=True)

    # arrow
    if num != "5":
        add_text(slide, "→", x + Inches(2.5), Inches(4.25), Inches(0.4), Inches(0.5),
                 size=22, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    x += Inches(2.55)

add_rect(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), DARK)
add_text(slide, "After enrollment: device receives APNs push within seconds for any remote action",
         Inches(0.5), Inches(7.03), Inches(12), Inches(0.28),
         size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Device Management
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Device Management", "Full remote control over enrolled macOS devices")

actions = [
    (Inches(0.3),  Inches(1.3), "🔒  Lock Device",    BLUE,
     ["Send PIN + message", "Immediate via APNs", "Device shows lock screen", "Audit logged"]),
    (Inches(3.5),  Inches(1.3), "🗑️  Erase Device",   ORANGE,
     ["Full wipe + reinstall", "Configurable PIN", "CheckOutWhenRemoved", "Irreversible action"]),
    (Inches(6.7),  Inches(1.3), "🔄  Restart",        GREEN,
     ["Graceful restart", "Queued via APNs", "Useful after profile install", "Instant delivery"]),
    (Inches(9.9),  Inches(1.3), "🔍  Query Info",     PURPLE,
     ["Serial, model, OS", "FileVault, Firewall", "Supervised status", "Updates available"]),

    (Inches(0.3),  Inches(3.8), "👥  Device Users",   RGBColor(0x0F, 0x76, 0x6E),
     ["List local users", "Admin / standard role", "Secure Token status", "Last login time"]),
    (Inches(3.5),  Inches(3.8), "📦  Patch Management", RGBColor(0xD9, 0x77, 0x06),
     ["Installed app inventory", "Available OS updates", "Critical update count", "Remote install"]),
    (Inches(6.7),  Inches(3.8), "🛡️  Compliance",    RGBColor(0x7C, 0x3A, 0xED),
     ["Policy evaluation", "Pass / Fail / Unknown", "Per-device report", "Fleet dashboard"]),
    (Inches(9.9),  Inches(3.8), "📋  Audit Logs",     DARK,
     ["All admin actions", "Actor email + IP", "Resource + changes", "Export as CSV"]),
]

for l, t, title, color, items in actions:
    add_rect(slide, l, t, Inches(2.9), Inches(3.2), WHITE)
    add_rect(slide, l, t, Inches(2.9), Inches(0.45), color)
    add_text(slide, title, l + Inches(0.1), t + Inches(0.06), Inches(2.7), Inches(0.35),
             size=13, bold=True, color=WHITE)
    y2 = t + Inches(0.55)
    for item in items:
        add_text(slide, f"✓  {item}", l + Inches(0.12), y2, Inches(2.65), Inches(0.3),
                 size=12, color=RGBColor(0x1E, 0x29, 0x3B))
        y2 += Inches(0.58)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PSSO
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
header_bar(slide, "Platform SSO — Microsoft Entra ID",
           "Passwordless macOS login with corporate credentials")

# Left: what it is
add_rect(slide, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.8), RGBColor(0x1E, 0x29, 0x3B))
add_text(slide, "What is PSSO?", Inches(0.5), Inches(1.4), Inches(5.4), Inches(0.45),
         size=18, bold=True, color=BLUE)
psso_points = [
    "macOS 13+ native Single Sign-On framework",
    "Login to Mac using Microsoft Entra ID credentials",
    "Secure Enclave key — no password stored on disk",
    "Seamless SSO to all Microsoft 365 / Azure apps",
    "Requires Microsoft Company Portal installed",
    "Delivered via com.apple.extensiblesso MDM profile",
    "Registration token from Entra automates device join",
]
y = Inches(1.95)
for pt in psso_points:
    add_text(slide, f"→  {pt}", Inches(0.5), y, Inches(5.4), Inches(0.35),
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(0.45)

# Right: profile key fields
add_rect(slide, Inches(6.5), Inches(1.3), Inches(6.4), Inches(5.8), RGBColor(0x0F, 0x17, 0x2A))
add_text(slide, "Profile Payload (mobileconfig)", Inches(6.7), Inches(1.4), Inches(6.0), Inches(0.45),
         size=16, bold=True, color=BLUE)

code_lines = [
    ("PayloadType:", "com.apple.extensiblesso"),
    ("ExtensionIdentifier:", "com.microsoft.CompanyPortalMac"),
    ("TeamIdentifier:", "UBF8T346G9"),
    ("Type:", "Credential"),
    ("AuthenticationMethod:", "UserSecureEnclaveKey"),
    ("EnableCreateUserAtLogin:", "true"),
    ("RegistrationToken:", "<from Entra portal>"),
    ("AdministratorGroups:", '["IT Admins"]'),
]
y = Inches(1.95)
for key, val in code_lines:
    add_text(slide, key, Inches(6.7), y, Inches(3.2), Inches(0.32),
             size=12, bold=True, color=RGBColor(0x7D, 0xD3, 0xFC))
    add_text(slide, val, Inches(9.9), y, Inches(2.8), Inches(0.32),
             size=12, color=RGBColor(0xBB, 0xF7, 0xD0))
    y += Inches(0.42)

add_rect(slide, Inches(6.5), Inches(6.3), Inches(6.4), Inches(0.6), BLUE)
add_text(slide, "Push to all enrolled devices in one click from the Policies page",
         Inches(6.7), Inches(6.38), Inches(6.0), Inches(0.42),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Software Distribution
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Software Distribution", "Upload once, deploy to any device via self-service portal")

# Flow steps
steps = [
    (BLUE,   "1. Upload Package", ["Admin uploads .pkg/.dmg/.zip", "Up to 4 GB supported", "Stored securely on server", "Versioning & description"]),
    (PURPLE, "2. Portal Catalog",  ["Appears in self-service portal", "User sees 'IT Managed' apps", "One-click install request", "Token-authenticated access"]),
    (ORANGE, "3. Admin Approval",  ["Request appears in dashboard", "Admin reviews + approves", "Or set auto-approve", "Audit logged"]),
    (GREEN,  "4. Agent Installs",  ["Agent picks up job in 30s", "Runs installer silently", "Reports exit code + output", "Status: installing → done"]),
]

x = Inches(0.25)
for color, title, items in steps:
    add_rect(slide, x, Inches(1.3), Inches(3.0), Inches(5.5), WHITE)
    add_rect(slide, x, Inches(1.3), Inches(3.0), Inches(0.5), color)
    add_text(slide, title, x + Inches(0.12), Inches(1.35), Inches(2.75), Inches(0.4),
             size=14, bold=True, color=WHITE)
    y2 = Inches(1.92)
    for item in items:
        add_text(slide, f"✓  {item}", x + Inches(0.15), y2, Inches(2.7), Inches(0.35),
                 size=13, color=DARK)
        y2 += Inches(0.65)
    if color != GREEN:
        add_text(slide, "→", x + Inches(3.05), Inches(3.8), Inches(0.3), Inches(0.5),
                 size=24, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    x += Inches(3.3)

# Bottom: agent detail
add_rect(slide, Inches(0.25), Inches(7.0), Inches(12.8), Inches(0.38), DARK)
add_text(slide,
         "Management Agent — pure bash LaunchDaemon, no Python or Xcode required. "
         "Bootstrap: curl -sSLG -d auth=<token> <url> | sudo bash",
         Inches(0.4), Inches(7.04), Inches(12.4), Inches(0.3),
         size=11, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — USB Block Policy
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "USB Block Policy", "Block or allow external storage — fleet-wide or per device")

# Left column
add_rect(slide, Inches(0.3), Inches(1.3), Inches(5.9), Inches(5.7), WHITE)
add_text(slide, "How It Works", Inches(0.5), Inches(1.4), Inches(5.5), Inches(0.4),
         size=17, bold=True, color=BLUE)
how = [
    ("macOS 12–13", "com.apple.systemuiserver with mount-controls. Denies harddisk-external, disk-image, DVD, CD, BD."),
    ("macOS 14+",   "com.apple.security.diskaccess with DenyExternalStorage: true. Full enforcement on supervised devices."),
    ("Both payloads sent together", "Single profile push covers all macOS versions automatically."),
    ("Deterministic identifier", "com.mdmsaas.usb.block.profile.{tenant_id} — ensures RemoveProfile always finds the right profile."),
    ("Removal disallowed", "Profile cannot be removed by user — only via MDM RemoveProfile command."),
]
y = Inches(1.95)
for title, desc in how:
    add_text(slide, f"• {title}", Inches(0.5), y, Inches(5.5), Inches(0.3),
             size=13, bold=True, color=DARK)
    add_text(slide, f"  {desc}", Inches(0.5), y + Inches(0.3), Inches(5.5), Inches(0.38),
             size=12, color=MID_GRAY, wrap=True)
    y += Inches(0.88)

# Right column
add_rect(slide, Inches(6.6), Inches(1.3), Inches(6.4), Inches(2.5), DARK)
add_text(slide, "Fleet Push (All Devices)", Inches(6.8), Inches(1.4), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=ORANGE)
add_text(slide, "Policies page → USB Block → Push to All\nQueues InstallProfile for every enrolled device simultaneously.\nAPNs push wakes each device instantly.",
         Inches(6.8), Inches(1.88), Inches(6.0), Inches(0.85),
         size=13, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

add_rect(slide, Inches(6.6), Inches(4.0), Inches(6.4), Inches(2.5), RGBColor(0x1E, 0x29, 0x3B))
add_text(slide, "Per-Device Control", Inches(6.8), Inches(4.1), Inches(6.0), Inches(0.4),
         size=15, bold=True, color=GREEN)
add_text(slide, "Device Detail page → Block USB / Remove USB Block\n\n"
         "Block: InstallProfile queued for that device only.\n"
         "Remove: RemoveProfile with deterministic identifier — works reliably every time.",
         Inches(6.8), Inches(4.58), Inches(6.0), Inches(1.2),
         size=13, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

add_rect(slide, Inches(6.6), Inches(6.6), Inches(6.4), Inches(0.7), BLUE)
add_text(slide, "Requires supervised device for silent enforcement.\n"
         "Unsupervised: macOS shows user-dismissable warning dialog.",
         Inches(6.8), Inches(6.65), Inches(6.0), Inches(0.6),
         size=12, color=WHITE, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — JIT Admin Access
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
header_bar(slide, "JIT Admin Access", "Temporary, audited admin elevation — auto-revoked on expiry")

stages = [
    (BLUE,   "Request",    "User opens self-service portal on their Mac.\nSelects their account, enters a reason and duration (1–24h).\nRequest sent to admin dashboard."),
    (ORANGE, "Approve",    "Admin reviews request in dashboard.\nApproves with optional duration override.\nrevoke_at timestamp set in database."),
    (GREEN,  "Elevate",    "Management agent picks up dseditgroup job.\nAdds user to 'admin' group on macOS.\nUser is now local admin for the duration."),
    (PURPLE, "Auto-Revoke","Background worker checks every 60 seconds.\nWhen revoke_at ≤ now: removes from admin group.\nWebhook notification sent (Slack/Teams/Discord)."),
]

x = Inches(0.3)
for color, stage, desc in stages:
    add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(5.0), RGBColor(0x1E, 0x29, 0x3B))
    add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(0.55), color)
    add_text(slide, stage, x + Inches(0.12), Inches(1.35), Inches(2.65), Inches(0.42),
             size=18, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.12), Inches(2.0), Inches(2.65), Inches(4.0),
             size=13, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)
    if stage != "Auto-Revoke":
        add_text(slide, "→", x + Inches(3.0), Inches(3.6), Inches(0.3), Inches(0.5),
                 size=26, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    x += Inches(3.28)

add_rect(slide, Inches(0.3), Inches(6.55), Inches(12.7), Inches(0.75), RGBColor(0x1E, 0x29, 0x3B))
add_text(slide, "Security principles: Least privilege • Time-bound access • Full audit trail • No persistent admin accounts",
         Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.55),
         size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Compliance Engine
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Compliance Engine", "Policy-based compliance evaluation across your entire fleet")

# Policy rules
add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.2), Inches(5.7), WHITE)
add_text(slide, "Policy Rules", Inches(0.5), Inches(1.4), Inches(3.8), Inches(0.4),
         size=16, bold=True, color=BLUE)
rules = [
    "FileVault encryption enabled",
    "Firewall enabled",
    "Gatekeeper enforced",
    "PSSO registered",
    "Screen lock configured",
    "Max check-in age (hours)",
    "Critical updates threshold",
]
y = Inches(1.95)
for rule in rules:
    add_text(slide, f"☑  {rule}", Inches(0.5), y, Inches(3.8), Inches(0.35),
             size=13, color=DARK)
    y += Inches(0.6)

# Status
add_rect(slide, Inches(4.8), Inches(1.3), Inches(3.8), Inches(2.5), WHITE)
add_text(slide, "Device Status", Inches(5.0), Inches(1.4), Inches(3.4), Inches(0.4),
         size=16, bold=True, color=BLUE)
statuses = [
    ("✅ Compliant",      GREEN,  "All rules passing"),
    ("❌ Non-Compliant",  ORANGE, "One or more rules failing"),
    ("❓ Unknown",        MID_GRAY, "Not yet evaluated"),
]
y = Inches(1.9)
for label, color, desc in statuses:
    add_text(slide, label, Inches(5.0), y, Inches(2.0), Inches(0.35),
             size=13, bold=True, color=color)
    add_text(slide, desc, Inches(5.0), y + Inches(0.32), Inches(3.4), Inches(0.28),
             size=11, color=MID_GRAY)
    y += Inches(0.75)

# Fleet summary
add_rect(slide, Inches(4.8), Inches(4.05), Inches(3.8), Inches(2.9), DARK)
add_text(slide, "Fleet Dashboard", Inches(5.0), Inches(4.15), Inches(3.4), Inches(0.4),
         size=15, bold=True, color=WHITE)
fleet = ["Total enrolled devices", "# Compliant", "# Non-compliant", "# Unknown", "Active policies list"]
y = Inches(4.65)
for item in fleet:
    add_text(slide, f"→  {item}", Inches(5.0), y, Inches(3.4), Inches(0.32),
             size=12, color=RGBColor(0xCB, 0xD5, 0xE1))
    y += Inches(0.44)

# Frameworks
add_rect(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.7), WHITE)
add_text(slide, "Supported Frameworks", Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.4),
         size=15, bold=True, color=BLUE)
frameworks = [
    ("CIS macOS Benchmark",    "Center for Internet Security hardening"),
    ("NIST 800-53",            "Federal security controls"),
    ("PCI DSS v4",             "Payment card industry standard"),
    ("ISO 27001",              "Information security management"),
    ("Custom",                 "Define your own policy rules"),
]
y = Inches(1.95)
for fw, desc in frameworks:
    add_text(slide, fw, Inches(9.2), y, Inches(3.6), Inches(0.3),
             size=13, bold=True, color=DARK)
    add_text(slide, desc, Inches(9.2), y + Inches(0.28), Inches(3.6), Inches(0.28),
             size=11, color=MID_GRAY)
    y += Inches(0.82)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Multi-Tenancy & Security
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
header_bar(slide, "Multi-Tenancy & Security", "Built for SaaS — complete isolation between customers")

cards = [
    (Inches(0.3),  Inches(1.3), "Tenant Isolation", BLUE,
     ["Every DB query: WHERE tenant_id = :id", "tenant_id from JWT, never from URL params",
      "MDM endpoints: UDID → tenant_id lookup", "No cross-tenant data leakage possible"]),
    (Inches(4.5),  Inches(1.3), "Authentication", DARK,
     ["Dashboard: JWT Bearer tokens", "MDM endpoints: mTLS device certificates",
      "Agent: per-device agent tokens", "Portal: per-device tokens (no JWT)"]),
    (Inches(8.7),  Inches(1.3), "Profile Security", GREEN,
     ["Profiles signed with CMS PKCS#7", "SHA-256 signature, DER encoded",
      "macOS verifies signature on install", "CA cert embedded in signature bag"]),
    (Inches(0.3),  Inches(4.3), "Audit Trail", PURPLE,
     ["All admin actions logged", "Actor email + IP address", "Resource type + changes (JSONB)",
      "Filterable by action + resource type"]),
    (Inches(4.5),  Inches(4.3), "Secrets Management", ORANGE,
     ["APNs certs via AWS Secrets Manager", "JWT key in environment variable",
      "DB password never in code", "LocalStack for dev (no AWS needed)"]),
    (Inches(8.7),  Inches(4.3), "Network Security", RGBColor(0x0F, 0x76, 0x6E),
     ["HTTPS everywhere (Caddy / ALB)", "APNs: TLS 1.2+ to api.push.apple.com",
      "MDM commands: device cert verified", "No raw plist bodies in logs (PII)"]),
]

for l, t, title, color, items in cards:
    add_rect(slide, l, t, Inches(3.8), Inches(2.75), WHITE)
    add_rect(slide, l, t, Inches(3.8), Inches(0.45), color)
    add_text(slide, title, l + Inches(0.12), t + Inches(0.06), Inches(3.55), Inches(0.35),
             size=14, bold=True, color=WHITE)
    y2 = t + Inches(0.55)
    for item in items:
        add_text(slide, f"✓  {item}", l + Inches(0.12), y2, Inches(3.55), Inches(0.32),
                 size=12, color=DARK)
        y2 += Inches(0.48)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Demo Flow
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK)
header_bar(slide, "Live Demo Walkthrough", "End-to-end demonstration of key capabilities")

demo_steps = [
    (BLUE,   "Step 1",  "Enroll a Device",
     "Generate enrollment token → Open URL in Safari on Mac → Install profile → Device appears enrolled"),
    (PURPLE, "Step 2",  "Remote Actions",
     "Query device info → View serial/OS/model → Lock device with PIN → View command status in real time"),
    (ORANGE, "Step 3",  "Push PSSO Profile",
     "Policies → Platform SSO → Configure Entra tenant → Push to all devices → Verify profile on Mac"),
    (GREEN,  "Step 4",  "Software Request",
     "Install agent via bootstrap URL → Open self-service portal → Request Zoom/Firefox → Approve → Auto-installs"),
    (RGBColor(0x0F, 0x76, 0x6E), "Step 5", "USB Block",
     "Device page → Block USB → Plug in USB stick → Blocked! → Remove USB Block → USB works again"),
    (RGBColor(0x7C, 0x3A, 0xED), "Step 6", "JIT Admin Access",
     "Portal → Request admin access (2h) → Approve in dashboard → Elevated → Auto-revoked after timeout"),
]

x = Inches(0.3)
row = 0
for i, (color, step, title, desc) in enumerate(demo_steps):
    col = i % 3
    row = i // 3
    l = Inches(0.3) + col * Inches(4.3)
    t = Inches(1.3) + row * Inches(2.8)
    add_rect(slide, l, t, Inches(4.1), Inches(2.55), RGBColor(0x1E, 0x29, 0x3B))
    add_rect(slide, l, t, Inches(4.1), Inches(0.48), color)
    add_text(slide, f"{step}: {title}", l + Inches(0.12), t + Inches(0.06),
             Inches(3.85), Inches(0.38), size=14, bold=True, color=WHITE)
    add_text(slide, desc, l + Inches(0.12), t + Inches(0.6), Inches(3.85), Inches(1.7),
             size=12, color=RGBColor(0xCB, 0xD5, 0xE1), wrap=True)

add_rect(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.32), BLUE)
add_text(slide, "Total demo time: ~15 minutes  |  All features running live on real macOS device",
         Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.26),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
dark_slide(slide)
add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), BLUE)
add_rect(slide, 0, Inches(7.44), SLIDE_W, Inches(0.06), BLUE)

add_text(slide, "MDM SaaS", Inches(0.5), Inches(1.0), Inches(12), Inches(1.5),
         size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "One platform to manage, secure, and support every Mac in your organisation.",
         Inches(0.5), Inches(2.5), Inches(12), Inches(0.7),
         size=22, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

caps = [
    "Enrollment", "Remote Actions", "PSSO / Entra ID", "Software Distribution",
    "USB Policy", "Compliance", "JIT Admin Access", "Audit Logs",
]
y = Inches(3.5)
row_caps = [caps[:4], caps[4:]]
for row in row_caps:
    x2 = Inches(1.0)
    for cap in row:
        add_rect(slide, x2, y, Inches(2.6), Inches(0.5), BLUE)
        add_text(slide, cap, x2, y + Inches(0.06), Inches(2.6), Inches(0.38),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x2 += Inches(2.85)
    y += Inches(0.68)

add_text(slide, "Built with FastAPI · PostgreSQL · Next.js · Apple APNs · AWS ECS",
         Inches(0.5), Inches(5.6), Inches(12), Inches(0.4),
         size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "Thank You",
         Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
         size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
out = "/Users/manikandank/Downloads/mdm-saas/MDM_SaaS_Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
